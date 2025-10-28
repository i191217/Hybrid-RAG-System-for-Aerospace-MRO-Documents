#!/usr/bin/env python3
"""
Text extraction service for individual documents.
Extracts text extraction logic from pipeline for API use.
"""

import os
import sys
import logging
import io
from pathlib import Path
from typing import Optional
import numpy as np
from tqdm import tqdm

# Add current directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import configuration
from core.config import config

# Optional imports with fallbacks
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import easyocr
    EASYOCR_AVAILABLE = True
except Exception as e:
    EASYOCR_AVAILABLE = False
    logging.getLogger(__name__).error(f"Failed to import EasyOCR: {e!r}")

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from tqdm import tqdm
    TQDM_AVAILABLE = True
except ImportError:
    TQDM_AVAILABLE = False

class TextExtractionService:
    """Service for extracting text from various document formats."""
    
    def __init__(self):
        """Initialize the text extraction service."""
        self.logger = logging.getLogger(__name__)
        self.easyocr_reader = None
        
    async def initialize(self):
        """Initialize OCR components."""
        try:
            # Initialize EasyOCR reader if available
            if EASYOCR_AVAILABLE:
                try:
                    self.easyocr_reader = easyocr.Reader(['en'])
                    self.logger.info("EasyOCR reader initialized successfully")
                except Exception as e:
                    self.logger.warning(f"Failed to initialize EasyOCR reader: {e}")
            else:
                self.logger.warning("EasyOCR not available. PDF text extraction will use fallback methods.")
                
        except Exception as e:
            self.logger.error(f"Error initializing text extraction service: {e}")
            raise

    async def cleanup(self):
        """Cleanup resources."""
        # No specific cleanup needed for current implementation
        pass

    async def extract_text(self, file_path: Path, extension: str) -> str:
        """
        Extract text from a file based on its extension.
        
        Args:
            file_path: Path to the file
            extension: File extension (with dot, e.g., '.pdf')
            
        Returns:
            Extracted text content
        """
        self.logger.info(f"Extracting text from: {file_path.name}")
        self.logger.info(f"File type: {extension}")
        
        try:
            if extension == '.pdf':
                return await self._extract_from_pdf(file_path)
            elif extension in ['.xlsx', '.xls']:
                return await self._extract_from_excel(file_path)
            elif extension in ['.docx', '.doc']:
                return await self._extract_from_word(file_path)
            elif extension == '.txt':
                return await self._extract_from_text(file_path)
            elif extension == '.csv':
                return await self._extract_from_csv(file_path)
            elif extension == '.pptx':
                return await self._extract_from_powerpoint(file_path)
            else:
                self.logger.warning(f"Unsupported file type: {extension}")
                return f"Unsupported file type: {extension}"
                
        except Exception as e:
            self.logger.error(f"Failed to extract text from {file_path.name}: {e}")
            return f"Error extracting text from {file_path.name}: {str(e)}"

    async def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files using EasyOCR and PyMuPDF."""
        self.logger.info(f"Extracting PDF content using EasyOCR...")
        
        try:
            # Check if required libraries are available
            if not PYMUPDF_AVAILABLE:
                self.logger.error("PyMuPDF not available. Install: pip install PyMuPDF")
                return await self._extract_from_pdf_fallback(file_path)
            
            if not EASYOCR_AVAILABLE or not self.easyocr_reader:
                self.logger.warning("EasyOCR not available. Using fallback method...")
                return await self._extract_from_pdf_fallback(file_path)
            
            if not PIL_AVAILABLE:
                self.logger.error("PIL not available. Install: pip install Pillow")
                return await self._extract_from_pdf_fallback(file_path)
            
            # Convert PDF to images using PyMuPDF
            self.logger.info(f"Using PyMuPDF + EasyOCR...")
            images = await self._pdf_to_images_pymupdf(file_path)
            
            if not images:
                self.logger.warning("Failed to convert PDF to images. Using fallback method...")
                return await self._extract_from_pdf_fallback(file_path)
            
            # Extract text from each page using EasyOCR
            page_texts = []
            failed_pages = []
            
            # Add progress bar for OCR processing
            self.logger.info(f"Starting OCR processing for {len(images)} pages...")
            page_iterator = tqdm(enumerate(images), total=len(images), 
                               desc="OCR Processing", unit="page") if len(images) > 1 else enumerate(images)
            
            for page_num, image in page_iterator:
                try:
                    if len(images) > 1:
                        page_iterator.set_description(f"OCR Processing page {page_num + 1}/{len(images)}")
                    
                    self.logger.debug(f"Processing page {page_num + 1} with EasyOCR...")
                    
                    # Convert PIL image to numpy array for EasyOCR
                    img_array = np.array(image)
                    
                    # Perform OCR
                    results = self.easyocr_reader.readtext(img_array, detail=0)
                    
                    # Join all text results
                    text = ' '.join(results)
                    
                    if text.strip():
                        page_texts.append(f"=== PAGE {page_num + 1} ===\n{text}\n")
                        self.logger.debug(f"Page {page_num + 1}: {len(text)} characters extracted")
                        if len(images) > 1:
                            page_iterator.set_postfix({
                                'chars': len(text),
                                'success': len(page_texts) - len(failed_pages),
                                'failed': len(failed_pages)
                            })
                    else:
                        self.logger.warning(f"Page {page_num + 1}: No text extracted")
                        failed_pages.append(page_num + 1)
                        page_texts.append(f"=== PAGE {page_num + 1} ===\n[No text extracted]\n")
                
                except Exception as e:
                    self.logger.warning(f"Page {page_num + 1}: EasyOCR failed - {e}")
                    failed_pages.append(page_num + 1)
                    page_texts.append(f"=== PAGE {page_num + 1} ===\n[OCR extraction failed: {e}]\n")
            
            # Combine all page texts
            combined_text = '\n'.join(page_texts)
            
            self.logger.info(f"PDF extraction complete using EasyOCR:")
            self.logger.info(f"- Total pages: {len(images)}")
            self.logger.info(f"- Successfully processed: {len(images) - len(failed_pages)}")
            self.logger.info(f"- Failed pages: {len(failed_pages)}")
            self.logger.info(f"- Total characters: {len(combined_text)}")
            
            return combined_text
            
        except Exception as e:
            self.logger.error(f"EasyOCR PDF extraction failed: {e}")
            self.logger.info("Attempting fallback method...")
            return await self._extract_from_pdf_fallback(file_path)

    async def _pdf_to_images_pymupdf(self, pdf_path: Path, dpi: int = 300) -> list:
        """Convert PDF pages to images using PyMuPDF."""
        try:
            self.logger.debug(f"Converting PDF to images using PyMuPDF...")
            
            # Open PDF with PyMuPDF
            doc = fitz.open(str(pdf_path))
            images = []
            
            # Add progress bar for image conversion
            self.logger.info(f"Converting {len(doc)} pages to images...")
            page_range = tqdm(range(len(doc)), desc="Converting to images", unit="page") if len(doc) > 1 else range(len(doc))
            
            # Convert each page to image
            for page_num in page_range:
                try:
                    if len(doc) > 1:
                        page_range.set_description(f"Converting page {page_num + 1}/{len(doc)} to image")
                    
                    page = doc.load_page(page_num)
                    
                    # Create transformation matrix for DPI
                    mat = fitz.Matrix(dpi / 72, dpi / 72)
                    
                    # Render page as image
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Convert to PIL Image
                    img_data = pix.tobytes("ppm")
                    image = Image.open(io.BytesIO(img_data))
                    images.append(image)
                    
                    if len(doc) > 1:
                        page_range.set_postfix({'converted': len(images)})
                        
                except Exception as e:
                    self.logger.warning(f"Failed to convert page {page_num + 1}: {e}")
                    continue
            
            doc.close()
            self.logger.debug(f"Converted {len(images)} pages to images")
            return images
            
        except Exception as e:
            self.logger.error(f"PyMuPDF PDF to image conversion failed: {e}")
            return []

    async def _extract_from_pdf_fallback(self, file_path: Path) -> str:
        """Fallback PDF extraction using traditional methods."""
        self.logger.info(f"Using fallback PDF extraction methods...")
        
        try:
            # Try PyPDF2 first
            import PyPDF2
            self.logger.info(f"Using PyPDF2...")
            
            text_content = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                self.logger.info(f"PDF has {len(pdf_reader.pages)} pages")
                
                # Add progress bar for page processing
                page_iterator = tqdm(enumerate(pdf_reader.pages), total=len(pdf_reader.pages), 
                                   desc="Processing PDF pages", unit="page") if len(pdf_reader.pages) > 1 else enumerate(pdf_reader.pages)
                
                for page_num, page in page_iterator:
                    try:
                        if len(pdf_reader.pages) > 1:
                            page_iterator.set_description(f"Processing page {page_num + 1}/{len(pdf_reader.pages)}")
                        
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(page_text)
                            self.logger.debug(f"Page {page_num + 1}: {len(page_text)} chars")
                            if len(pdf_reader.pages) > 1:
                                page_iterator.set_postfix({'chars': len(page_text), 'pages_done': len(text_content)})
                        else:
                            self.logger.warning(f"Page {page_num + 1}: No text extracted")
                    except Exception as e:
                        self.logger.warning(f"Page {page_num + 1}: Error - {e}")
                        continue
            
            combined_text = '\n'.join(text_content)
            self.logger.info(f"PDF extraction complete: {len(combined_text)} characters from {len(text_content)} pages")
            return combined_text
            
        except ImportError:
            self.logger.warning("PyPDF2 not available, trying pdfplumber...")
            try:
                import pdfplumber
                self.logger.info(f"Using pdfplumber...")
                
                text_content = []
                with pdfplumber.open(file_path) as pdf:
                    self.logger.info(f"PDF has {len(pdf.pages)} pages")
                    
                    # Add progress bar for page processing
                    page_iterator = tqdm(enumerate(pdf.pages), total=len(pdf.pages), 
                                       desc="Processing PDF pages (pdfplumber)", unit="page") if len(pdf.pages) > 1 else enumerate(pdf.pages)
                    
                    for page_num, page in page_iterator:
                        try:
                            if len(pdf.pages) > 1:
                                page_iterator.set_description(f"Processing page {page_num + 1}/{len(pdf.pages)}")
                            
                            page_text = page.extract_text()
                            if page_text and page_text.strip():
                                text_content.append(page_text)
                                self.logger.debug(f"Page {page_num + 1}: {len(page_text)} chars")
                                if len(pdf.pages) > 1:
                                    page_iterator.set_postfix({'chars': len(page_text), 'pages_done': len(text_content)})
                            else:
                                self.logger.warning(f"Page {page_num + 1}: No text extracted")
                        except Exception as e:
                            self.logger.warning(f"Page {page_num + 1}: Error - {e}")
                            continue
                
                combined_text = '\n'.join(text_content)
                self.logger.info(f"PDF extraction complete: {len(combined_text)} characters from {len(text_content)} pages")
                return combined_text
                
            except ImportError:
                self.logger.error("No PDF libraries available. Install: pip install PyPDF2 or pip install pdfplumber")
                return f"PDF extraction failed: No PDF libraries available for {file_path.name}"
        
        except Exception as e:
            self.logger.error(f"PDF extraction failed: {e}")
            return f"PDF extraction failed: {str(e)}"

    async def _extract_from_excel(self, file_path: Path) -> str:
        """Extract text from Excel files."""
        self.logger.info(f"Extracting Excel content...")
        
        try:
            import openpyxl
            self.logger.info(f"Using openpyxl...")
            
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            self.logger.info(f"Excel has {len(workbook.sheetnames)} sheets: {workbook.sheetnames}")
            
            all_text = []
            
            # Add progress bar for sheet processing
            sheet_iterator = tqdm(workbook.sheetnames, desc="Processing Excel sheets", unit="sheet") if len(workbook.sheetnames) > 1 else workbook.sheetnames
            
            for sheet_name in sheet_iterator:
                if len(workbook.sheetnames) > 1:
                    sheet_iterator.set_description(f"Processing sheet: {sheet_name}")
                
                self.logger.info(f"Processing sheet: {sheet_name}")
                sheet = workbook[sheet_name]
                
                sheet_text = []
                sheet_text.append(f"=== SHEET: {sheet_name} ===")
                
                # Get the used range
                if sheet.max_row > 0 and sheet.max_column > 0:
                    self.logger.info(f"Range: {sheet.max_row} rows x {sheet.max_column} columns")
                    
                    # Add progress bar for row processing if sheet is large
                    row_iterator = tqdm(sheet.iter_rows(values_only=True), total=sheet.max_row, 
                                       desc=f"Processing rows in {sheet_name}", unit="row") if sheet.max_row > 100 else sheet.iter_rows(values_only=True)
                    
                    for row in row_iterator:
                        row_text = []
                        for cell_value in row:
                            if cell_value is not None:
                                row_text.append(str(cell_value))
                        
                        if row_text:  # Only add non-empty rows
                            sheet_text.append('\t'.join(row_text))
                    
                    self.logger.info(f"Extracted {len(sheet_text)-1} non-empty rows")
                    if len(workbook.sheetnames) > 1:
                        sheet_iterator.set_postfix({'rows': len(sheet_text)-1})
                else:
                    self.logger.warning(f"Sheet {sheet_name} appears to be empty")
                
                if len(sheet_text) > 1:  # More than just the header
                    all_text.extend(sheet_text)
                    all_text.append("")  # Add spacing between sheets
            
            combined_text = '\n'.join(all_text)
            self.logger.info(f"Excel extraction complete: {len(combined_text)} characters")
            return combined_text
            
        except ImportError:
            self.logger.warning("openpyxl not available, trying pandas...")
            try:
                import pandas as pd
                self.logger.info(f"Using pandas...")
                
                # Read all sheets
                excel_file = pd.ExcelFile(file_path)
                self.logger.info(f"Excel has {len(excel_file.sheet_names)} sheets: {excel_file.sheet_names}")
                
                all_text = []
                
                # Add progress bar for sheet processing
                sheet_iterator = tqdm(excel_file.sheet_names, desc="Processing Excel sheets (pandas)", unit="sheet") if len(excel_file.sheet_names) > 1 else excel_file.sheet_names
                
                for sheet_name in sheet_iterator:
                    if len(excel_file.sheet_names) > 1:
                        sheet_iterator.set_description(f"Processing sheet: {sheet_name}")
                    
                    self.logger.info(f"Processing sheet: {sheet_name}")
                    try:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        
                        if not df.empty:
                            sheet_text = [f"=== SHEET: {sheet_name} ==="]
                            
                            # Add column headers
                            headers = [str(col) for col in df.columns if str(col) != 'nan']
                            if headers:
                                sheet_text.append('\t'.join(headers))
                            
                            # Add data rows with progress bar for large sheets
                            row_iterator = tqdm(df.iterrows(), total=len(df), 
                                               desc=f"Processing rows in {sheet_name}", unit="row") if len(df) > 100 else df.iterrows()
                            
                            for _, row in row_iterator:
                                row_values = [str(val) for val in row.values if str(val) != 'nan' and val is not None]
                                if row_values:
                                    sheet_text.append('\t'.join(row_values))
                            
                            all_text.extend(sheet_text)
                            all_text.append("")  # Add spacing
                            self.logger.info(f"Extracted {len(df)} rows x {len(df.columns)} columns")
                            if len(excel_file.sheet_names) > 1:
                                sheet_iterator.set_postfix({'rows': len(df), 'cols': len(df.columns)})
                        else:
                            self.logger.warning(f"Sheet {sheet_name} is empty")
                    
                    except Exception as e:
                        self.logger.warning(f"Failed to read sheet {sheet_name}: {e}")
                        continue
                
                combined_text = '\n'.join(all_text)
                self.logger.info(f"Excel extraction complete: {len(combined_text)} characters")
                return combined_text
                
            except ImportError:
                self.logger.error("No Excel libraries available. Install: pip install openpyxl or pip install pandas")
                return f"Excel extraction failed: No Excel libraries available for {file_path.name}"
        
        except Exception as e:
            self.logger.error(f"Excel extraction failed: {e}")
            return f"Excel extraction failed: {str(e)}"

    async def _extract_from_word(self, file_path: Path) -> str:
        """Extract text from Word documents."""
        self.logger.info(f"Extracting Word document content...")
        
        try:
            import docx
            self.logger.info(f"Using python-docx...")
            
            doc = docx.Document(file_path)
            
            text_content = []
            
            # Extract paragraphs
            self.logger.info(f"Document has {len(doc.paragraphs)} paragraphs")
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
                    
            # Extract tables
            if doc.tables:
                self.logger.info(f"Document has {len(doc.tables)} tables")
                for table_num, table in enumerate(doc.tables):
                    text_content.append(f"\n=== TABLE {table_num + 1} ===")
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_content.append('\t'.join(row_text))
            
            combined_text = '\n'.join(text_content)
            self.logger.info(f"Word extraction complete: {len(combined_text)} characters")
            return combined_text
            
        except ImportError:
            self.logger.error("python-docx not available. Install: pip install python-docx")
            return f"Word extraction failed: python-docx library not available for {file_path.name}"
        
        except Exception as e:
            self.logger.error(f"Word extraction failed: {e}")
            return f"Word extraction failed: {str(e)}"

    async def _extract_from_text(self, file_path: Path) -> str:
        """Extract text from plain text files."""
        self.logger.info(f"Reading text file...")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                self.logger.info(f"Text file read: {len(content)} characters")
                return content
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin1') as file:
                    content = file.read()
                    self.logger.info(f"Text file read (latin1): {len(content)} characters")
                    return content
            except Exception as e:
                self.logger.error(f"Failed to read text file: {e}")
                return f"Text extraction failed: {str(e)}"

    async def _extract_from_csv(self, file_path: Path) -> str:
        """Extract text from CSV files."""
        self.logger.info(f"Reading CSV file...")
        
        try:
            import pandas as pd
            self.logger.info(f"Using pandas for CSV...")
            
            df = pd.read_csv(file_path)
            
            text_content = [f"=== CSV: {file_path.name} ==="]
            
            # Add column headers
            headers = [str(col) for col in df.columns]
            text_content.append('\t'.join(headers))
            
            # Add data rows
            for _, row in df.iterrows():
                row_values = [str(val) for val in row.values if str(val) != 'nan' and val is not None]
                if row_values:
                    text_content.append('\t'.join(row_values))
            
            combined_text = '\n'.join(text_content)
            self.logger.info(f"CSV extraction complete: {len(combined_text)} characters from {len(df)} rows")
            return combined_text
            
        except ImportError:
            # Fallback to basic CSV reading
            self.logger.info("pandas not available, using basic CSV reading...")
            try:
                import csv
                
                text_content = [f"=== CSV: {file_path.name} ==="]
                
                with open(file_path, 'r', encoding='utf-8') as file:
                    csv_reader = csv.reader(file)
                    for row in csv_reader:
                        if row:  # Skip empty rows
                            text_content.append('\t'.join(row))
                
                combined_text = '\n'.join(text_content)
                self.logger.info(f"CSV extraction complete: {len(combined_text)} characters")
                return combined_text
                
            except Exception as e:
                self.logger.error(f"CSV extraction failed: {e}")
                return f"CSV extraction failed: {str(e)}"
        
        except Exception as e:
            self.logger.error(f"CSV extraction failed: {e}")
            return f"CSV extraction failed: {str(e)}"

    async def _extract_from_powerpoint(self, file_path: Path) -> str:
        """Extract text from PowerPoint files."""
        self.logger.info(f"Extracting PowerPoint content...")
        
        try:
            from pptx import Presentation
            self.logger.info(f"Using python-pptx...")
            
            prs = Presentation(file_path)
            
            text_content = [f"=== POWERPOINT: {file_path.name} ==="]
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"\n=== SLIDE {slide_num + 1} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the header
                    text_content.extend(slide_text)
            
            combined_text = '\n'.join(text_content)
            self.logger.info(f"PowerPoint extraction complete: {len(combined_text)} characters from {len(prs.slides)} slides")
            return combined_text
            
        except ImportError:
            self.logger.error("python-pptx not available. Install: pip install python-pptx")
            return f"PowerPoint extraction failed: python-pptx library not available for {file_path.name}"
        
        except Exception as e:
            self.logger.error(f"PowerPoint extraction failed: {e}")
            return f"PowerPoint extraction failed: {str(e)}" 